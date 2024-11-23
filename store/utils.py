from PyPDF2 import PdfReader
from pptx import Presentation

def search_in_pdf_file(pdf_path, searched_data):
    """
    Searches for specific data in a PDF file.

    Args:
        pdf_path (str): The path to the PDF file.
        searched_data (str): The data to search for in the PDF file.

    Returns:
        bool: True if the searched data is found in the PDF file, False otherwise.
    """
    pdf_file_obj = open(pdf_path, 'rb')
    pdf_reader = PdfReader(pdf_file_obj)
    num_pages = len(pdf_reader.pages)
    count = 0
    text = ""
    # The while loop will read each page.
    while count < num_pages:
        page_obj = pdf_reader.pages[count]
        count += 1
        text += page_obj.extract_text()
        if searched_data in text.lower():
            return True
    return False

def search_in_pptx_file(pdf_path,searched_data):
    """
    Searches for specific data in a PowerPoint presentation.

    Args:
        pdf_path (str): The path to the PowerPoint presentation file.
        searched_data (str): The data to search for in the PowerPoint presentation.

    Returns:
        bool: True if the searched data is found in the PowerPoint presentation, False otherwise.
    """
    prs = Presentation(pdf_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and searched_data in shape.text.lower():
                return True
    return False

